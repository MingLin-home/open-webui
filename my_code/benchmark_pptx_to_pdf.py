import os
import sys
import requests
import base64
import time
from pptx import Presentation

def has_image_in_pptx(pptx_filename):
    if not pptx_filename.lower().endswith('.pptx'):
        raise ValueError("The input file must be a .pptx file.")
    prs = Presentation(pptx_filename)

    # Iterate over all slides and their shapes to check for images
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE = 13
                return True
    return False


def pptx_to_pdf_base64(pptx_path):
    url = 'http://127.0.0.1:2004/request'
    with open(pptx_path, 'rb') as f:
        files = {'file': (os.path.basename(pptx_path), f, 'application/vnd.openxmlformats-officedocument.presentationml.presentation')}
        data = {'convert-to': 'pdf'}
        response = requests.post(url, files=files, data=data)
        response.raise_for_status()
        pdf_bytes = response.content
        b64_pdf = base64.b64encode(pdf_bytes).decode('utf-8')
        return b64_pdf


def save_base64_as_pdf(base64_str, pdf_path):
    pdf_bytes = base64.b64decode(base64_str)
    with open(pdf_path, 'wb') as f:
        f.write(pdf_bytes)


if __name__ == "__main__":
    input_fn = '/Users/mmilin/tmp/Test Open-WebUI KnowledgeBase/311introductiontomachinelearning.pptx'

    time_start = time.perf_counter()
    has_image = has_image_in_pptx(pptx_filename=input_fn)
    time_end = time.perf_counter()
    elapsed_time_has_image = time_end - time_start

    time_start = time.perf_counter()
    pptx_to_pdf_base64(pptx_path=input_fn)
    time_end = time.perf_counter()
    elapsed_time_pptx_to_pdf_base64 = time_end - time_start

    print(f"elapsed_time_has_image={elapsed_time_has_image}, elapsed_time_pptx_to_pdf_base64={elapsed_time_pptx_to_pdf_base64}")